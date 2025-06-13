# app.py
import os
import io
from flask import Flask, request, jsonify, send_file, render_template_string
from pptx import Presentation
import base64
from werkzeug.utils import secure_filename

# --- Flask App Initialization ---
app = Flask(__name__)

# --- Configuration ---
# Define the folder to store uploaded PowerPoint files
UPLOAD_FOLDER = "uploads"
# Define allowed file extensions
ALLOWED_EXTENSIONS = {"pptx"}

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# --- Create Upload Directory ---
# Create the upload folder if it doesn't already exist
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)


# --- Helper Function ---
def allowed_file(filename):
    """Checks if the uploaded file has a .pptx extension"""
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def replace_text_in_shape(shape, replacements):
    """Replaces placeholders in a single shape, including across multiple runs."""
    if not shape.has_text_frame:
        print("NOT REPLACE IN SHAPEEEEE")
        return

    full_text = shape.text_frame.text
    for placeholder, value in replacements.items():
        formatted_placeholder = f"{{{placeholder}}}"
        if formatted_placeholder in full_text:
            full_text = full_text.replace(formatted_placeholder, str(value))
    
    # Replace the entire text (this will overwrite individual runs, but ensures accuracy)
    shape.text_frame.clear()  # Clear existing runs
    shape.text_frame.text = full_text

def replace_text_in_table(shape, replacements):
    """Replaces placeholders within a table shape."""
    if not shape.has_table:
        print("NOT REPLACE IN TABLEEE")
        return

    # Iterate through each cell in the table
    for row in shape.table.rows:
        for cell in row.cells:
            # Each cell has a text_frame, similar to a regular shape
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, value in replacements.items():
                        formatted_placeholder = f"{{{placeholder}}}"
                        if formatted_placeholder in run.text:
                            run.text = run.text.replace(
                                formatted_placeholder, str(value)
                            )


# --- HTML Template for Upload Form ---
# A simple web page to test the upload functionality
HTML_TEMPLATE = """
<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1, shrink-to-fit=no">
  <title>PPTX Modifier API</title>
  <link href="https://cdn.jsdelivr.net/npm/tailwindcss@2.2.19/dist/tailwind.min.css" rel="stylesheet">
</head>
<body class="bg-gray-100 flex items-center justify-center h-screen">
  <div class="bg-white p-8 rounded-lg shadow-md w-full max-w-md">
    <h1 class="text-2xl font-bold mb-4 text-center text-gray-700">Upload PowerPoint File</h1>
    <p class="text-gray-600 mb-6 text-center">Upload a .pptx file to get started.</p>
    <form action="/upload" method="post" enctype="multipart/form-data">
      <div class="mb-4">
        <label for="file" class="block mb-2 text-sm font-medium text-gray-900">Choose file</label>
        <input type="file" name="file" id="file" class="block w-full text-sm text-gray-900 bg-gray-50 rounded-lg border border-gray-300 cursor-pointer focus:outline-none">
      </div>
      <button type="submit" class="w-full text-white bg-blue-600 hover:bg-blue-700 focus:ring-4 focus:ring-blue-300 font-medium rounded-lg text-sm px-5 py-2.5 text-center">
        Upload
      </button>
    </form>
  </div>
</body>
</html>
"""

# --- API Endpoints ---


@app.route("/")
def index():
    """Renders the HTML upload form."""
    return render_template_string(HTML_TEMPLATE)


@app.route("/upload", methods=["POST"])
def upload_file():
    """
    Endpoint: Handles the upload of a .pptx file.
    The uploaded file will be renamed and saved as 'TEMPLATE_SALES.pptx',
    replacing any existing file with the same name.
    """
    if "file" not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    # 1. This checks if the file is a valid type (e.g., .pptx)
    #    regardless of its original name.
    if file and allowed_file(file.filename):
        
        # 2. Here, we explicitly set the server-side filename. This is the
        #    "rename" step. Any uploaded file, be it "report.pptx" or
        #    "draft_v2.pptx", will be prepared to be saved as this name.
        filename = "TEMPLATE_SALES.pptx"
        
        save_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        
        # 3. The file.save() function saves the uploaded file to the 'save_path'.
        #    If a file at that path already exists, it is automatically
        #    overwritten. This is the "replace" step.
        file.save(save_path)

        return (
            jsonify(
                {
                    "message": "Template file uploaded and updated successfully.",
                    "filename": filename
                }
            ),
            200, # Using 200 OK is common for updates.
        )
    else:
        return jsonify({"error": "Invalid file type. Please upload a .pptx file."}), 400

@app.route("/process", methods=["POST"])
def process_file():
    filename = "TEMPLATE_SALES.pptx"
    """
    Endpoint 2: Processes the specified .pptx file.
    Replaces placeholders like {companyName} with values from the JSON body.
    """
    # Get the replacement data from the request body
    try:
        replacements = request.get_json()
        if not replacements:
            raise ValueError()
    except:
        return jsonify({"error": "Invalid or missing JSON body"}), 400

    # Construct the path to the target file
    filepath = os.path.join(app.config["UPLOAD_FOLDER"], secure_filename(filename))

    # Check if the file exists
    if not os.path.exists(filepath):
        return jsonify({"error": "File not found"}), 404

    try:
        # Open the presentation
        prs = Presentation(filepath)

        # Iterate through each slide and each shape to find and replace text
        for slide in prs.slides:
            for shape in slide.shapes:
                print("REPLACEMENTTT ",replacements)
                # Replace text in standard shapes and text boxes
                replace_text_in_shape(shape, replacements)
                # Replace text inside tables
                replace_text_in_table(shape, replacements)

        # Save the modified presentation to an in-memory stream
        file_stream = io.BytesIO()
        prs.save(file_stream)
        # Rewind the stream to the beginning
        file_stream.seek(0)

        # Encode the file to base64 and decode to a utf-8 string
        encoded_string = base64.b64encode(file_stream.read()).decode("utf-8")

        # Return a JSON response containing the base64 data.
        # This format is more robust for integration with workflow tools like n8n.
        return jsonify(
            {
                "message": "File processed successfully",
                "filename": f"modified_{filename}",
                "data": encoded_string,  # The base64 data for the file
                "mimetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
        )

    except Exception as e:
        # Catch potential errors from the python-pptx library or other issues
        return (
            jsonify(
                {"error": f"An error occurred while processing the file: {str(e)}"}
            ),
            500,
        )


# --- Run the App ---
if __name__ == "__main__":
    # Make sure to run in debug mode only for development
    app.run(host="0.0.0.0", port=8000, debug=True)
